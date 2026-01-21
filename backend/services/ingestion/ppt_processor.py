import os
from pptx import Presentation

class PPTProcessor:
    def __init__(self, config):
        self.config = config

    def process(self, file_path):
        """
        Processes a PPT file: extracts text per slide and renders slide images as JPGs.
        """
        import aspose.slides as slides
        
        # 1. Text extraction via python-pptx (cleaner)
        prs = Presentation(file_path)
        base_filename = os.path.basename(file_path).replace('.pptx', '').replace('.ppt', '')
        
        # 2. Rendering via aspose-slides
        with slides.Presentation(file_path) as aspose_prs:
            results = []
            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                text_content = self._extract_text_from_slide(slide)
                
                # Render slide to image
                image_name = f"{base_filename}_slide_{slide_number}.jpg"
                image_path = os.path.join(self.config.SLIDES_FOLDER, image_name)
                
                # Using aspose to save as JPG
                aspose_slide = aspose_prs.slides[i]
                image = aspose_slide.get_image(1.0, 1.0)
                try:
                    image.save(image_path, slides.ImageFormat.JPEG)
                    if not os.path.exists(image_path):
                        print(f"  Warning: Slide image {image_name} failed to save at {image_path}")
                except Exception as e:
                    print(f"  Error saving slide image {image_name}: {e}")
                
                # Calculate relative URL for frontend (serving via /api/slides/)
                image_url = f"/api/slides/{image_name}"
                
                results.append({
                    "page_number": slide_number,
                    "text_content": text_content,
                    "image_url": image_url,
                    "slide_layout": slide.slide_layout.name if slide.slide_layout else "unknown"
                })
            
            return results

    def _extract_text_from_slide(self, slide):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
        return "\n".join(text_runs)
